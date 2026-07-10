"""
presentation.py

Generates an executive PowerPoint from all JSON reports found in:

    reports/weekly/*.json

Usage:
    python src/presentation.py

Output:
    reports/executive_project_health_report.pptx
"""

import json
from collections import Counter
from pathlib import Path
from statistics import mean

from pptx import Presentation
from pptx.util import Inches


REPORTS_DIR = Path("reports/weekly")
OUTPUT_FILE = Path("reports/executive_project_health_report.pptx")


# ==========================================================
# LOAD REPORTS
# ==========================================================

def load_reports():

    reports = []

    for file in REPORTS_DIR.glob("*.json"):

        try:

            with open(file, "r", encoding="utf-8") as f:
                reports.append(json.load(f))

        except Exception as e:

            print(f"Skipping {file}: {e}")

    return reports


# ==========================================================
# HELPERS
# ==========================================================

def rag_counts(reports):

    counts = Counter()

    for r in reports:
        counts[r["official_rag"]] += 1

    return counts


def portfolio_health(counts):

    red = counts.get("Red", 0)
    amber = counts.get("Amber", 0)

    if red > 0:
        return "High Risk"

    if amber > 0:
        return "Moderate Risk"

    return "Healthy"


def average_confidence(reports):

    vals = []

    for r in reports:

        val = (
            r.get("scorer_detail", {})
             .get("confidence_score")
        )

        if val is not None:
            vals.append(val)

    return round(mean(vals), 2) if vals else None


def average_data_quality(reports):

    vals = []

    for r in reports:

        val = (
            r.get("scorer_detail", {})
             .get("data_quality_score")
        )

        if val is not None:
            vals.append(val)

    return round(mean(vals), 1) if vals else None


# ==========================================================
# TREND ANALYSIS
# ==========================================================

def derive_trends(reports):

    trends = []

    overdue_projects = 0
    variance_projects = 0
    comment_risk_projects = 0
    critical_path_projects = 0

    for r in reports:

        features = r.get("features", {})

        overdue = (
            features
            .get("overdue_task_ratio", {})
            .get("value")
        )

        variance = (
            features
            .get("variance_health", {})
            .get("value")
        )

        comment_risk = (
            features
            .get("comment_risk_signal", {})
            .get("value")
        )

        milestone = (
            features
            .get("milestone_health", {})
            .get("value")
        )

        if overdue is not None and overdue > 0.10:
            overdue_projects += 1

        if variance is not None and variance > 7:
            variance_projects += 1

        if comment_risk is not None and comment_risk >= 0.5:
            comment_risk_projects += 1

        if milestone is not None and milestone > 0:
            critical_path_projects += 1

    total = len(reports)

    if overdue_projects > 0:
        trends.append(
            f"Overdue task accumulation observed in "
            f"{overdue_projects}/{total} projects."
        )

    if variance_projects > 0:
        trends.append(
            f"Schedule variance is elevated in "
            f"{variance_projects}/{total} projects."
        )

    if comment_risk_projects > 0:
        trends.append(
            f"Stakeholder comments indicate delivery concerns in "
            f"{comment_risk_projects}/{total} projects."
        )

    if critical_path_projects > 0:
        trends.append(
            f"Critical-path delays are present in "
            f"{critical_path_projects}/{total} projects."
        )

    avg_quality = average_data_quality(reports)

    if avg_quality is not None and avg_quality >= 75:
        trends.append(
            "Portfolio data quality remains sufficiently high "
            "for decision support."
        )

    return trends


# ==========================================================
# RISKS
# ==========================================================

def derive_portfolio_risks(reports):

    risks = Counter()

    for r in reports:

        negatives = (
            r.get("insights", {})
             .get("negative_drivers", [])
        )

        for item in negatives:
            risks[item] += 1

    return [
        x[0]
        for x in risks.most_common(5)
    ]


# ==========================================================
# RECOMMENDATIONS
# ==========================================================

def derive_recommendations(reports):

    recs = []

    critical_projects = []

    blocker_projects = []

    for r in reports:

        name = r["project_name"]

        milestone = (
            r["features"]
            .get("milestone_health", {})
            .get("value")
        )

        blocker = (
            r["features"]
            .get("blocker_density", {})
            .get("value")
        )

        if milestone is not None and milestone > 0:
            critical_projects.append(name)

        if blocker is not None and blocker > 0.20:
            blocker_projects.append(name)

    if critical_projects:

        recs.append(
            "Review recovery plans for projects with "
            "overdue critical-path activities."
        )

    if blocker_projects:

        recs.append(
            "Investigate dependency bottlenecks and "
            "critical-path blocker accumulation."
        )

    recs.append(
        "Continue monitoring overdue task growth "
        "before milestone impact occurs."
    )

    recs.append(
        "Maintain current governance practices on "
        "healthy projects."
    )

    return recs


# ==========================================================
# PPT HELPERS
# ==========================================================

def add_title_slide(prs, reports):

    slide = prs.slides.add_slide(
        prs.slide_layouts[0]
    )

    counts = rag_counts(reports)

    slide.shapes.title.text = (
        "Project Health Reporting Agent"
    )

    slide.placeholders[1].text = (
        f"Projects Assessed: {len(reports)}\n"
        f"Green: {counts.get('Green',0)}\n"
        f"Amber: {counts.get('Amber',0)}\n"
        f"Red: {counts.get('Red',0)}\n\n"
        f"Portfolio Health: "
        f"{portfolio_health(counts)}"
    )


def add_bullets_slide(prs, title, bullets):

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()

    if not bullets:
        bullets = ["No notable items identified."]

    for i, item in enumerate(bullets):

        if i == 0:
            tf.text = item
        else:
            tf.add_paragraph().text = item


def add_attention_projects_slide(prs, reports):

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = (
        "Projects Requiring Attention"
    )

    tf = slide.placeholders[1].text_frame
    tf.clear()

    attention = [
        r for r in reports
        if r["official_rag"] != "Green"
    ]

    if not attention:

        tf.text = (
            "No projects currently require "
            "special intervention."
        )

        return

    first = True

    for r in attention:

        lines = []

        lines.append(
            f"{r['project_name']} "
            f"({r['official_rag']})"
        )

        negatives = (
            r["insights"]
            .get("negative_drivers", [])
        )

        lines.extend(
            negatives[:3]
        )

        text = "\n".join(lines)

        if first:
            tf.text = text
            first = False
        else:
            tf.add_paragraph().text = text


def add_dashboard_slide(prs, reports):

    slide = prs.slides.add_slide(
        prs.slide_layouts[5]
    )

    slide.shapes.title.text = (
        "Portfolio Dashboard"
    )

    rows = len(reports) + 1
    cols = 4

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.5),
        Inches(1.2),
        Inches(8.5),
        Inches(2)
    ).table

    headers = [
        "Project",
        "RAG",
        "Confidence",
        "Data Quality"
    ]

    for c, h in enumerate(headers):
        table.cell(0, c).text = h

    for i, r in enumerate(reports, start=1):

        table.cell(
            i, 0
        ).text = r["project_name"]

        table.cell(
            i, 1
        ).text = r["official_rag"]

        table.cell(
            i, 2
        ).text = str(
            round(
                r["scorer_detail"]
                ["confidence_score"] * 100
            )
        ) + "%"

        table.cell(
            i, 3
        ).text = str(
            r["scorer_detail"]
            ["data_quality_score"]
        )


def add_methodology_slide(prs, reports):

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = (
        "Methodology & Confidence"
    )

    avg_conf = average_confidence(reports)
    avg_quality = average_data_quality(reports)

    slide.placeholders[1].text = (
        "Workflow\n"
        "Project Plan → Feature Engineering → "
        "Comment Analysis → Deterministic RAG "
        "Engine → Executive Reporting\n\n"
        f"Portfolio Confidence: "
        f"{round(avg_conf*100)}%\n"
        f"Portfolio Data Quality: "
        f"{avg_quality}/100\n\n"
        "RAG classification is determined "
        "algorithmically.\n"
        "LLMs provide explanation only."
    )


# ==========================================================
# MAIN
# ==========================================================

def build_presentation():

    reports = load_reports()

    if not reports:
        raise RuntimeError(
            "No JSON reports found in "
            "reports/weekly/"
        )

    prs = Presentation()

    add_title_slide(prs, reports)

    add_bullets_slide(
        prs,
        "Cross-Project Trends",
        derive_trends(reports)
    )

    add_bullets_slide(
        prs,
        "Emerging Portfolio Risks",
        derive_portfolio_risks(reports)
    )

    add_attention_projects_slide(
        prs,
        reports
    )

    add_bullets_slide(
        prs,
        "Executive Recommendations",
        derive_recommendations(reports)
    )

    add_dashboard_slide(
        prs,
        reports
    )

    add_methodology_slide(
        prs,
        reports
    )

    OUTPUT_FILE.parent.mkdir(
        parents=True,
        exist_ok=True
    )

    prs.save(OUTPUT_FILE)

    print(
        f"Presentation written to: "
        f"{OUTPUT_FILE}"
    )


if __name__ == "__main__":
    build_presentation()

